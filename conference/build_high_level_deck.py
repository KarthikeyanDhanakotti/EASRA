"""Build the EASRA High-Level Architecture slide deck for LinkedIn / conferences.

Idempotent. Re-run whenever the palette, layer names, or footer change.

Output: conference/EASRA-High-Level-Architecture.pptx
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PALETTE = {
    # (fill, border, text)
    "entry":    ("#EAF3FF", "#3B7DDD", "#0B3D91"),
    "neutral":  ("#ECF0F1", "#2C3E50", "#1B2A3A"),
    "model":    ("#F5EAF7", "#8E44AD", "#4A1E5C"),
    "data":     ("#E8F8F5", "#16A085", "#0B5C4B"),
    "security": ("#FDECEA", "#C0392B", "#6B1A11"),
    "obs":      ("#FFF6E5", "#B58900", "#5A4200"),
    "cache":    ("#FFF6E5", "#B58900", "#5A4200"),
    "biz":      ("#EAF3FF", "#3B7DDD", "#0B3D91"),
    "titlebar": ("#2C3E50", "#2C3E50", "#FFFFFF"),
    "footer":   ("#F5F7FA", "#DDE2E8", "#5A6470"),
    "hero_bg":  ("#0B1F3A", "#0B1F3A", "#FFFFFF"),
    "hero_acc": ("#3B7DDD", "#3B7DDD", "#FFFFFF"),
}


def hex_to_rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def add_box(slide, left, top, width, height, fill, border, text,
            title=None, font_size=11, title_size=13, bold_title=True,
            border_weight=1.25, rounded=True, text_color=None,
            title_color=None, align_left=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill)
    shp.line.color.rgb = hex_to_rgb(border)
    shp.line.width = Pt(border_weight)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    tf.text = ""
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if align_left else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = title
        r.font.size = Pt(title_size)
        r.font.bold = bold_title
        r.font.name = "Segoe UI"
        r.font.color.rgb = hex_to_rgb(title_color or text_color or "#1B2A3A")
        if text:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT if align_left else PP_ALIGN.CENTER
            r2 = p2.add_run()
            r2.text = text
            r2.font.size = Pt(font_size)
            r2.font.name = "Segoe UI"
            r2.font.color.rgb = hex_to_rgb(text_color or "#1B2A3A")
    else:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if align_left else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text or ""
        r.font.size = Pt(font_size)
        r.font.name = "Segoe UI"
        r.font.color.rgb = hex_to_rgb(text_color or "#1B2A3A")
    return shp


def add_text(slide, left, top, width, height, text, size=12, bold=False,
             color="#1B2A3A", align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             italic=False, name="Segoe UI"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = name
    r.font.color.rgb = hex_to_rgb(color)
    return tb


# ---------- slide builders ---------- #

def build_hero(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # background band
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(PALETTE["hero_bg"][0])
    bg.line.fill.background()

    # accent bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.55), prs.slide_width, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(PALETTE["hero_acc"][0])
    bar.line.fill.background()

    # eyebrow
    add_text(
        slide, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.4),
        "OPEN ARCHITECTURE STANDARD  ·  v0.1.0",
        size=13, bold=True, color="#3B7DDD", align=PP_ALIGN.LEFT,
    )

    # main title
    add_text(
        slide, Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.8),
        "EASRA", size=88, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(2.95), Inches(11.5), Inches(1.0),
        "Enterprise AI Systems Reference Architecture",
        size=36, bold=False, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )

    # tagline
    add_text(
        slide, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.65),
        "A vendor-neutral open standard for production-grade Enterprise AI.",
        size=20, italic=True, color="#B8CDEA", align=PP_ALIGN.LEFT,
    )

    # summary badges
    badges = [
        ("16", "Layers"),
        ("4", "Trust Boundaries"),
        ("3", "Cross-cutting Planes"),
        ("10", "Deliverables"),
    ]
    x = Inches(0.9)
    y = Inches(5.0)
    bw = Inches(2.6)
    bh = Inches(1.15)
    gap = Inches(0.2)
    for value, label in badges:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb("#12335F")
        shp.line.color.rgb = hex_to_rgb("#3B7DDD")
        shp.line.width = Pt(1.25)
        tf = shp.text_frame
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        r1.font.size = Pt(30)
        r1.font.bold = True
        r1.font.color.rgb = hex_to_rgb("#FFFFFF")
        r1.font.name = "Segoe UI"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(12)
        r2.font.color.rgb = hex_to_rgb("#B8CDEA")
        r2.font.name = "Segoe UI"
        x += bw + gap

    # author
    add_text(
        slide, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.35),
        "Karthikeyan Dhanakotti  ·  github.com/KarthikeyanDhanakotti/EASRA  "
        "·  CC-BY-4.0 (docs) + Apache-2.0 (code)",
        size=11, color="#8FA9CE", align=PP_ALIGN.LEFT,
    )


def build_architecture(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    W = prs.slide_width
    H = prs.slide_height

    # white background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
    bg.line.fill.background()

    # title bar
    title_h = Inches(0.65)
    title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, title_h)
    title.fill.solid()
    title.fill.fore_color.rgb = hex_to_rgb(PALETTE["titlebar"][0])
    title.line.fill.background()
    add_text(
        slide, Inches(0.35), Inches(0.05), Inches(12.6), Inches(0.55),
        "EASRA  ·  High-Level Architecture",
        size=20, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.35), Inches(0.05), W - Inches(0.7), Inches(0.55),
        "16 layers  ·  4 trust boundaries  ·  3 cross-cutting planes",
        size=11, color="#B8CDEA", align=PP_ALIGN.RIGHT,
    )

    # user pill
    top = title_h + Inches(0.15)
    user_h = Inches(0.45)
    add_box(
        slide,
        Inches(0.35), top, Inches(12.6), user_h,
        PALETTE["entry"][0], PALETTE["entry"][1],
        "Enterprise Users  ·  Web  ·  Mobile  ·  API  ·  Chat  ·  Voice",
        font_size=12, text_color=PALETTE["entry"][2],
    )

    # main content area
    area_top = top + user_h + Inches(0.15)
    area_bottom = H - Inches(1.05)  # leave room for platform strip + footer
    area_height = area_bottom - area_top
    area_left = Inches(0.35)
    area_right = W - Inches(0.35)

    # split: request flow (left ~62%) + cross-cutting (right ~38%)
    gap = Inches(0.18)
    left_w = int((area_right - area_left - gap) * 0.62)
    right_w = (area_right - area_left - gap) - left_w

    # ----- left column: request flow layers -----
    flow_layers = [
        ("L0", "Channels & User Experience",
         "Web · Mobile · API · Chat · Voice · SDK", "entry"),
        ("L1", "Edge · Gateway · Identity",
         "CDN · WAF · Global LB · API Gateway · OIDC/OAuth2 · Rate Limit · Session · Request Validator",
         "entry"),
        ("L2", "AI Orchestration",
         "Router · Single Agent · Multi-Agent Coordinator · Planner / Workflow Engine",
         "neutral"),
        ("L3", "Prompt Intelligence",
         "Context Builder · Prompt Builder · Versioned Prompt Registry",
         "neutral"),
        # L4 + L5 side by side
        ("split", "L4 Memory & Context  |  L5 Knowledge & Retrieval",
         "Short/Long/Session/Semantic Memory  ·  User Profile  ||  "
         "Retrieval Router  ·  Vector · BM25 · SQL · Graph  ·  Reranker",
         "data"),
        ("L6", "AI Models & Model Router",
         "Model Router (capability · cost · latency · policy)  ·  Foundation + Small/Specialised Models  ·  Model Registry",
         "model"),
        ("L7", "Tooling & Actions  (MCP)",
         "Tool Router · MCP Client ⇄ Servers · Enterprise APIs · Impact-Class Enforcer · Tool Registry",
         "data"),
        ("L8", "Guardrails & Safety",
         "Input · Prompt · Tool-Arg · Output guardrails",
         "security"),
        ("L9", "Verification  &  Streaming",
         "Grounding · Citation · Factuality · Format · Policy · Safety  →  Response Formatter · Streaming Engine",
         "security"),
    ]

    n = len(flow_layers)
    row_gap = Inches(0.08)
    total_gap = row_gap * (n - 1)
    row_h = (area_height - total_gap) / n

    for i, (code, title_text, body, kind) in enumerate(flow_layers):
        y = area_top + (row_h + row_gap) * i
        fill, border, text_col = PALETTE[kind]
        # code chip
        chip_w = Inches(0.72)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, area_left, y, chip_w, row_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        chip.line.width = Pt(0.5)
        ctf = chip.text_frame
        ctf.margin_left = Emu(0)
        ctf.margin_right = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code if code != "split" else "L4·L5"
        cr.font.size = Pt(13)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        # body box
        body_left = area_left + chip_w + Inches(0.06)
        body_w = left_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, body,
            title=title_text,
            title_size=12, font_size=10,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # ----- right column: cross-cutting planes -----
    right_left = area_left + left_w + gap
    cc_layers = [
        ("L10", "Observability", "obs",
         "OpenTelemetry Collector  ·  Traces · Metrics · Logs  ·  Tokens · Latency · Cost  ·  Agent · Prompt · Tool trace  ·  Evaluation · Alerts"),
        ("L13", "Security  &  Zero Trust", "security",
         "AuthN · AuthZ  ·  Secrets · Encryption  ·  PII Detection  ·  Policy Engine  ·  Audit  ·  Compliance"),
        ("L14", "Governance · Risk · Compliance", "neutral",
         "Policy Engine  ·  Model · Prompt · Tool registries  ·  NIST AI RMF · ISO 42001 · EU AI Act · OWASP LLM Top 10"),
        ("L15", "Business Outcomes  &  Value", "biz",
         "KPIs  ·  Value attribution  ·  Cost / Benefit  ·  Adoption"),
    ]

    # slight header for the right column
    hdr_h = Inches(0.32)
    add_text(
        slide, right_left, area_top - Inches(0.02), right_w, hdr_h,
        "CROSS-CUTTING PLANES  ·  attach to every layer",
        size=11, bold=True, color="#5A6470",
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    cc_area_top = area_top + Inches(0.30)
    cc_area_h = area_height - Inches(0.30)
    cc_gap = Inches(0.10)
    cc_n = len(cc_layers)
    cc_row_h = (cc_area_h - cc_gap * (cc_n - 1)) / cc_n

    for i, (code, title_text, kind, body) in enumerate(cc_layers):
        y = cc_area_top + (cc_row_h + cc_gap) * i
        fill, border, text_col = PALETTE[kind]
        chip_w = Inches(0.72)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, right_left, y, chip_w, cc_row_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        chip.line.width = Pt(0.5)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code
        cr.font.size = Pt(13)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = right_left + chip_w + Inches(0.06)
        body_w = right_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, cc_row_h,
            fill, border, body,
            title=title_text,
            title_size=12, font_size=10,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # ----- bottom strip: platform substrate -----
    strip_top = area_bottom + Inches(0.08)
    strip_h = Inches(0.62)
    strip_w_each = (area_right - area_left - gap * 2) / 3

    substrate = [
        ("L11  Performance · Caching · Cost",
         "Prompt · Semantic · Embedding · Memory · Response caches  ·  Cost Ledger",
         "cache"),
        ("L12  LLMOps & Delivery",
         "Git → CI (lint · sec · prompt-tests · eval) → Artefact Registry → CD (canary · blue/green · shadow · rollback)",
         "data"),
        ("Infrastructure Substrate",
         "Azure · AWS · GCP · On-Prem · Hybrid  ·  Kubernetes  ·  GPU / NPU / CPU  ·  Kafka · EventGrid · PubSub  ·  Object Storage  ·  Multi-region · DR",
         "neutral"),
    ]
    for i, (t, body, kind) in enumerate(substrate):
        x = area_left + (strip_w_each + gap) * i
        fill, border, text_col = PALETTE[kind]
        add_box(
            slide, x, strip_top, strip_w_each, strip_h,
            fill, border, body,
            title=t, title_size=11, font_size=9,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # footer
    add_text(
        slide, Inches(0.35), H - Inches(0.32), W - Inches(0.7), Inches(0.28),
        "github.com/KarthikeyanDhanakotti/EASRA  ·  v0.1.0  ·  CC-BY-4.0 (docs) + Apache-2.0 (code)  ·  Karthikeyan Dhanakotti",
        size=9, color="#8A94A0", align=PP_ALIGN.LEFT,
    )


def build_layer_index(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
    bg.line.fill.background()

    # title bar
    title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.65))
    title.fill.solid()
    title.fill.fore_color.rgb = hex_to_rgb(PALETTE["titlebar"][0])
    title.line.fill.background()
    add_text(
        slide, Inches(0.35), Inches(0.05), Inches(12.6), Inches(0.55),
        "EASRA  ·  The Sixteen Layers",
        size=20, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )

    layers = [
        ("L0", "Channels & User Experience", "entry",
         "Web, mobile, API, chat, voice, SDK — every way a user or system talks to the platform."),
        ("L1", "Edge · Gateway · Identity", "entry",
         "CDN, WAF, global load balancer, API gateway, OIDC/OAuth2, rate limiting, session, request validation."),
        ("L2", "AI Orchestration", "neutral",
         "Router, single agent, multi-agent coordinator, planner/workflow engine."),
        ("L3", "Prompt Intelligence", "neutral",
         "Context builder, prompt builder, versioned prompt registry."),
        ("L4", "Memory & Context", "data",
         "Short-term, long-term, session, semantic memory and user profile."),
        ("L5", "Knowledge & Retrieval", "data",
         "Retrieval router across vector, BM25, SQL, graph; reranker; hybrid retrieval."),
        ("L6", "AI Models & Model Router", "model",
         "Capability/cost/latency/policy-aware routing over foundation + small/specialised models."),
        ("L7", "Tooling & Actions (MCP)", "data",
         "Tool router, MCP clients/servers, enterprise APIs, impact-class enforcer, tool registry."),
        ("L8", "Guardrails & Safety", "security",
         "Input, prompt, tool-argument and output guardrails."),
        ("L9", "Verification", "security",
         "First-class verification layer: grounding, citation, factuality, format, policy, safety."),
        ("L10", "Observability", "obs",
         "OpenTelemetry-based traces, metrics, logs; tokens, latency, cost; agent/prompt/tool trace; evaluation, alerts."),
        ("L11", "Performance · Caching · Cost", "cache",
         "Prompt, semantic, embedding, memory, response caches; cost ledger."),
        ("L12", "LLMOps & Delivery", "data",
         "Git → CI (lint, sec-scan, prompt-tests, eval) → artefact registry → CD (canary, blue/green, shadow, rollback)."),
        ("L13", "Security & Zero Trust", "security",
         "AuthN/AuthZ, secrets, encryption, PII detection, policy engine, audit, compliance."),
        ("L14", "Governance · Risk · Compliance", "neutral",
         "Policy engine, model/prompt/tool registries, NIST AI RMF, ISO 42001, EU AI Act, OWASP LLM Top 10."),
        ("L15", "Business Outcomes & Value", "biz",
         "KPIs, value attribution, cost/benefit, adoption tracking."),
    ]

    # two-column grid: 8 rows × 2 cols
    top = Inches(0.9)
    left = Inches(0.35)
    col_gap = Inches(0.2)
    row_gap = Inches(0.08)
    cols = 2
    rows = 8
    total_w = W - Inches(0.7)
    col_w = (total_w - col_gap) / cols
    total_h = H - top - Inches(0.5)
    row_h = (total_h - row_gap * (rows - 1)) / rows

    for i, (code, name, kind, desc) in enumerate(layers):
        col = i // rows
        row = i % rows
        x = left + (col_w + col_gap) * col
        y = top + (row_h + row_gap) * row
        fill, border, text_col = PALETTE[kind]

        chip_w = Inches(0.72)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, chip_w, row_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code
        cr.font.size = Pt(13)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = x + chip_w + Inches(0.06)
        body_w = col_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, desc,
            title=name, title_size=11, font_size=9.5,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # footer
    add_text(
        slide, Inches(0.35), H - Inches(0.32), W - Inches(0.7), Inches(0.28),
        "Full spec: github.com/KarthikeyanDhanakotti/EASRA/tree/main/specification  ·  v0.1.0  ·  Karthikeyan Dhanakotti",
        size=9, color="#8A94A0", align=PP_ALIGN.LEFT,
    )


def build_square_architecture(prs: Presentation):
    """Single-slide 1:1 (10in x 10in) variant, optimised for LinkedIn mobile feed."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = prs.slide_width
    H = prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")
    bg.line.fill.background()

    # title bar
    title_h = Inches(0.9)
    title = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, title_h)
    title.fill.solid()
    title.fill.fore_color.rgb = hex_to_rgb(PALETTE["titlebar"][0])
    title.line.fill.background()
    add_text(
        slide, Inches(0.35), Inches(0.10), W - Inches(0.7), Inches(0.40),
        "EASRA",
        size=22, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.35), Inches(0.48), W - Inches(0.7), Inches(0.35),
        "Enterprise AI Systems Reference Architecture",
        size=14, color="#B8CDEA", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.35), Inches(0.10), W - Inches(0.7), Inches(0.80),
        "16 layers  ·  4 trust boundaries  ·  3 planes",
        size=11, color="#B8CDEA", align=PP_ALIGN.RIGHT,
    )

    # user pill
    top = title_h + Inches(0.18)
    user_h = Inches(0.45)
    add_box(
        slide,
        Inches(0.35), top, W - Inches(0.7), user_h,
        PALETTE["entry"][0], PALETTE["entry"][1],
        "Enterprise Users  ·  Web  ·  Mobile  ·  API  ·  Chat  ·  Voice",
        font_size=11, text_color=PALETTE["entry"][2],
    )

    area_top = top + user_h + Inches(0.15)
    area_bottom = H - Inches(1.15)  # room for substrate strip + footer
    area_height = area_bottom - area_top
    area_left = Inches(0.35)
    area_right = W - Inches(0.35)

    gap = Inches(0.14)
    # square canvas: use 60/40 split so cross-cutting labels still fit
    left_w = int((area_right - area_left - gap) * 0.60)
    right_w = (area_right - area_left - gap) - left_w

    # ---- request flow ---- #
    flow_layers = [
        ("L0", "Channels & UX", "Web · Mobile · API · Chat · Voice · SDK", "entry"),
        ("L1", "Edge · Gateway · Identity",
         "CDN · WAF · LB · Gateway · OIDC · Rate limit · Session · Validator",
         "entry"),
        ("L2", "AI Orchestration",
         "Router · Single agent · Multi-agent · Planner / Workflow",
         "neutral"),
        ("L3", "Prompt Intelligence",
         "Context builder · Prompt builder · Prompt registry",
         "neutral"),
        ("split", "L4 Memory  |  L5 Knowledge & Retrieval",
         "Short/Long/Session/Semantic memory  ||  Vector · BM25 · SQL · Graph · Reranker",
         "data"),
        ("L6", "AI Models & Model Router",
         "Router (capability · cost · latency · policy) · Foundation + specialised · Registry",
         "model"),
        ("L7", "Tooling & Actions (MCP)",
         "Tool router · MCP · Enterprise APIs · Impact-class enforcer · Registry",
         "data"),
        ("L8", "Guardrails & Safety",
         "Input · Prompt · Tool-arg · Output guardrails",
         "security"),
        ("L9", "Verification & Streaming",
         "Grounding · Citation · Factuality · Format · Policy · Safety  →  Formatter · Stream",
         "security"),
    ]

    n = len(flow_layers)
    row_gap = Inches(0.07)
    row_h = (area_height - row_gap * (n - 1)) / n

    for i, (code, title_text, body, kind) in enumerate(flow_layers):
        y = area_top + (row_h + row_gap) * i
        fill, border, text_col = PALETTE[kind]
        chip_w = Inches(0.72)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, area_left, y, chip_w, row_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        chip.line.width = Pt(0.5)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code if code != "split" else "L4·L5"
        cr.font.size = Pt(12)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = area_left + chip_w + Inches(0.06)
        body_w = left_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, body,
            title=title_text,
            title_size=11, font_size=9,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # ---- cross-cutting ---- #
    right_left = area_left + left_w + gap
    cc_layers = [
        ("L10", "Observability", "obs",
         "OTel  ·  Traces · Metrics · Logs  ·  Tokens · Latency · Cost  ·  Agent/Prompt/Tool trace  ·  Alerts"),
        ("L13", "Security & Zero Trust", "security",
         "AuthN · AuthZ  ·  Secrets · Encryption  ·  PII  ·  Policy engine  ·  Audit  ·  Compliance"),
        ("L14", "Governance · Risk · Compliance", "neutral",
         "Policy engine · Registries · NIST AI RMF · ISO 42001 · EU AI Act · OWASP LLM"),
        ("L15", "Business Outcomes & Value", "biz",
         "KPIs · Value attribution · Cost/benefit · Adoption"),
    ]

    add_text(
        slide, right_left, area_top - Inches(0.03), right_w, Inches(0.28),
        "CROSS-CUTTING PLANES",
        size=10, bold=True, color="#5A6470",
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
    )

    cc_area_top = area_top + Inches(0.27)
    cc_area_h = area_height - Inches(0.27)
    cc_gap = Inches(0.10)
    cc_n = len(cc_layers)
    cc_row_h = (cc_area_h - cc_gap * (cc_n - 1)) / cc_n

    for i, (code, title_text, kind, body) in enumerate(cc_layers):
        y = cc_area_top + (cc_row_h + cc_gap) * i
        fill, border, text_col = PALETTE[kind]
        chip_w = Inches(0.72)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, right_left, y, chip_w, cc_row_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        chip.line.width = Pt(0.5)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code
        cr.font.size = Pt(12)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = right_left + chip_w + Inches(0.06)
        body_w = right_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, cc_row_h,
            fill, border, body,
            title=title_text,
            title_size=11, font_size=9,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # ---- platform substrate strip ---- #
    strip_top = area_bottom + Inches(0.10)
    strip_h = Inches(0.75)
    strip_w_each = (area_right - area_left - gap * 2) / 3
    substrate = [
        ("L11 · Performance · Caching · Cost",
         "Prompt · Semantic · Embedding · Memory · Response caches  ·  Cost ledger",
         "cache"),
        ("L12 · LLMOps & Delivery",
         "Git → CI (lint · sec · prompt · eval) → Artefact Registry → CD (canary · shadow · rollback)",
         "data"),
        ("Infrastructure Substrate",
         "Azure · AWS · GCP · Hybrid  ·  Kubernetes  ·  GPU / NPU  ·  Kafka  ·  Multi-region · DR",
         "neutral"),
    ]
    for i, (t, body, kind) in enumerate(substrate):
        x = area_left + (strip_w_each + gap) * i
        fill, border, text_col = PALETTE[kind]
        add_box(
            slide, x, strip_top, strip_w_each, strip_h,
            fill, border, body,
            title=t, title_size=10, font_size=8.5,
            title_color=text_col, text_color=text_col,
            align_left=True,
        )

    # footer
    add_text(
        slide, Inches(0.35), H - Inches(0.30), W - Inches(0.7), Inches(0.26),
        "github.com/KarthikeyanDhanakotti/EASRA  ·  v0.1.0  ·  Karthikeyan Dhanakotti",
        size=9, color="#8A94A0", align=PP_ALIGN.LEFT,
    )


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_hero(prs)
    build_architecture(prs)
    build_layer_index(prs)

    out = Path(__file__).parent / "EASRA-High-Level-Architecture.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB)")

    # ----- square (1:1) variant for LinkedIn mobile feed ----- #
    prs_sq = Presentation()
    prs_sq.slide_width = Inches(10)
    prs_sq.slide_height = Inches(10)
    build_square_architecture(prs_sq)
    out_sq = Path(__file__).parent / "EASRA-High-Level-Architecture-Square.pptx"
    prs_sq.save(out_sq)
    print(f"Wrote {out_sq}  ({out_sq.stat().st_size / 1024:.1f} KB)")


if __name__ == "__main__":
    main()
