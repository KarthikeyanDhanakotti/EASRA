"""Build the EASRA LinkedIn carousel: cover, principles, layers, all key diagrams, CTA.

Idempotent. Re-run to regenerate `conference/EASRA-LinkedIn-Carousel.pptx`;
export to PDF via PowerPoint COM (see companion PowerShell one-liner).
"""
from __future__ import annotations
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from build_high_level_deck import (  # noqa: E402  reused primitives
    PALETTE,
    hex_to_rgb,
    add_box,
    add_text,
    build_hero,
    build_architecture,
    build_layer_index,
)


# ---------- helpers ---------- #

def slide_bg(prs: Presentation, color: str = "#FFFFFF"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(color)
    bg.line.fill.background()
    return slide, W, H


def title_bar(slide, W, subtitle_right: str | None = None,
              text: str = "EASRA  ·  Enterprise AI Systems Reference Architecture",
              height: float = 0.65):
    h = Inches(height)
    tb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, h)
    tb.fill.solid()
    tb.fill.fore_color.rgb = hex_to_rgb(PALETTE["titlebar"][0])
    tb.line.fill.background()
    add_text(
        slide, Inches(0.35), Inches(0.05), W - Inches(0.7), h - Inches(0.1),
        text, size=18, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    if subtitle_right:
        add_text(
            slide, Inches(0.35), Inches(0.05), W - Inches(0.7), h - Inches(0.1),
            subtitle_right, size=11, color="#B8CDEA", align=PP_ALIGN.RIGHT,
        )
    return h


def slide_footer(slide, W, H, text: str | None = None):
    footer = (
        text
        or "github.com/KarthikeyanDhanakotti/EASRA  ·  v0.1.0  ·  "
           "CC-BY-4.0 (docs) + Apache-2.0 (code)  ·  Karthikeyan Dhanakotti"
    )
    add_text(
        slide, Inches(0.35), H - Inches(0.32), W - Inches(0.7), Inches(0.28),
        footer, size=9, color="#8A94A0", align=PP_ALIGN.LEFT,
    )


def arrow_between(slide, x1, y1, x2, y2, color="#5A6470", weight=1.75):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2)
    )
    line.line.color.rgb = hex_to_rgb(color)
    line.line.width = Pt(weight)
    return line


# ---------- individual slide builders ---------- #

def build_problem(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Why EASRA?")

    add_text(
        slide, Inches(0.6), Inches(1.0), W - Inches(1.2), Inches(1.0),
        "Enterprise AI stacks are one-offs.",
        size=40, bold=True, color="#0B1F3A", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.6), Inches(2.0), W - Inches(1.2), Inches(0.8),
        "Every team is redesigning the same architecture — badly, in isolation, "
        "with different vocabulary.",
        size=18, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    problems = [
        ("No shared vocabulary",
         "Teams argue about ‘agent’ vs ‘workflow’ vs ‘orchestrator’ before writing a line of code."),
        ("Security bolted on",
         "Prompt injection, tool safety and PII treated as add-ons, not architectural concerns."),
        ("Verification confused with evaluation",
         "Offline metrics ship to prod. No first-class verification of grounding, citation, factuality."),
        ("Cost & observability afterthoughts",
         "Token spend, latency SLOs and tool traces retrofitted after the outage."),
    ]

    col_w = (W - Inches(1.4)) / 2
    row_h = Inches(1.5)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_y = Inches(3.2)
    for i, (t, d) in enumerate(problems):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + gap_x) * col
        y = start_y + (row_h + gap_y) * row
        add_box(
            slide, x, y, col_w, row_h,
            PALETTE["security"][0], PALETTE["security"][1], d,
            title=t, title_size=14, font_size=11,
            title_color=PALETTE["security"][2],
            text_color="#4A2C29", align_left=True,
        )

    slide_footer(slide, W, H)


def build_what_is(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="What is EASRA?")

    add_text(
        slide, Inches(0.6), Inches(1.0), W - Inches(1.2), Inches(1.0),
        "An open architecture standard for Enterprise AI.",
        size=36, bold=True, color="#0B1F3A", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.6), Inches(2.0), W - Inches(1.2), Inches(0.5),
        "Vendor-neutral. Production-grade. Standards-aligned.",
        size=16, italic=True, color="#3B7DDD", align=PP_ALIGN.LEFT,
    )

    quads = [
        ("Vendor-neutral",
         "Every layer maps cleanly to Azure, AWS, GCP and open-source. No lock-in.",
         "entry"),
        ("Security by design",
         "Zero trust, four trust boundaries, prompt-injection resistance in the layers, not on top.",
         "security"),
        ("Verification-first",
         "Grounding, citation, factuality, policy and safety verification separated from evaluation.",
         "obs"),
        ("Cost + observability by default",
         "Tokens, latency, cost, prompts, tools traced through OpenTelemetry from every component.",
         "data"),
    ]
    col_w = (W - Inches(1.4)) / 2
    row_h = Inches(1.55)
    start_y = Inches(2.85)
    for i, (t, d, kind) in enumerate(quads):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + (col_w + Inches(0.2)) * col
        y = start_y + (row_h + Inches(0.2)) * row
        fill, border, tcol = PALETTE[kind]
        add_box(
            slide, x, y, col_w, row_h,
            fill, border, d,
            title=t, title_size=14, font_size=11,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H)


def build_principles(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Ten design principles")

    principles = [
        ("P1", "Vendor Neutrality",
         "Logical layers; each maps to ≥ 2 independent implementations."),
        ("P2", "Security by Design",
         "Zero-trust, least-privilege, prompt-injection resistance are architectural."),
        ("P3", "Verification by Design",
         "Every AI output is verifiable; verification ≠ evaluation."),
        ("P4", "Observability by Default",
         "Traces, metrics, logs, token/cost accounting emitted from every component."),
        ("P5", "Loose Coupling",
         "Layers communicate only through published interfaces."),
        ("P6", "Externalized State",
         "Memory, sessions, caches live outside compute; compute is stateless."),
        ("P7", "Human-in-the-Loop",
         "Every autonomous action has escalation, override and audit paths."),
        ("P8", "Failure Isolation",
         "Layer failure degrades gracefully; no single layer cascades a full outage."),
        ("P9", "Cost Awareness",
         "Token, compute, storage cost are explicit design inputs."),
        ("P10", "Evolvability",
         "New models, tools, agents, channels added without redesigning the architecture."),
    ]

    cols = 2
    rows = 5
    top = Inches(0.95)
    left = Inches(0.35)
    col_gap = Inches(0.2)
    row_gap = Inches(0.08)
    total_w = W - Inches(0.7)
    col_w = (total_w - col_gap) / cols
    total_h = H - top - Inches(0.5)
    row_h = (total_h - row_gap * (rows - 1)) / rows

    for i, (code, name, desc) in enumerate(principles):
        col = i // rows
        row = i % rows
        x = left + (col_w + col_gap) * col
        y = top + (row_h + row_gap) * row
        fill, border, tcol = PALETTE["neutral"]

        chip_w = Inches(0.65)
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, chip_w, row_h)
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb("#3B7DDD")
        chip.line.color.rgb = hex_to_rgb("#3B7DDD")
        chip.line.width = Pt(0.5)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = code
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = x + chip_w + Inches(0.06)
        body_w = col_w - chip_w - Inches(0.06)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, desc,
            title=name, title_size=13, font_size=11,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H)


def build_trust_boundaries(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Four trust boundaries")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Every request crosses four boundaries. Each has its own attackers, controls, and failure modes.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    # Request path row
    path_y = Inches(1.6)
    path_h = Inches(0.7)
    boxes = [
        ("User", "#EAF3FF", "#3B7DDD", "#0B3D91"),
        ("L0 Channel", "#EAF3FF", "#3B7DDD", "#0B3D91"),
        ("L1 Gateway", "#EAF3FF", "#3B7DDD", "#0B3D91"),
        ("L2 Orch.", "#ECF0F1", "#2C3E50", "#1B2A3A"),
        ("L3 Prompt", "#ECF0F1", "#2C3E50", "#1B2A3A"),
        ("L6 Model Router", "#F5EAF7", "#8E44AD", "#4A1E5C"),
        ("L7 Tool Router", "#E8F8F5", "#16A085", "#0B5C4B"),
        ("L8 Guard", "#FDECEA", "#C0392B", "#6B1A11"),
        ("L9 Verify", "#FDECEA", "#C0392B", "#6B1A11"),
    ]
    total_w = W - Inches(0.7)
    n = len(boxes)
    gap = Inches(0.08)
    box_w = (total_w - gap * (n - 1)) / n
    for i, (label, fill, border, tcol) in enumerate(boxes):
        x = Inches(0.35) + (box_w + gap) * i
        add_box(
            slide, x, path_y, box_w, path_h, fill, border, label,
            font_size=10, text_color=tcol, title_size=10, bold_title=True,
        )

    # Trust boundary bands under the request path
    tbs = [
        ("TB-A · Edge Boundary",
         "untrusted → trusted", "#FDECEA", "#C0392B", "#6B1A11",
         "AuthN / AuthZ · WAF · rate limit · request validation · session"),
        ("TB-C · Model Boundary",
         "internal → provider", "#F5EAF7", "#8E44AD", "#4A1E5C",
         "Model registry · egress policy · prompt-inject defence · content filters · cost cap"),
        ("TB-D · Action Boundary",
         "reasoning → real world", "#E8F8F5", "#16A085", "#0B5C4B",
         "Tool registry · impact class · argument guardrails · HITL approval · audit"),
        ("TB-B · Response Boundary",
         "internal → observable", "#FDECEA", "#C0392B", "#6B1A11",
         "Output guardrails · verification verdict · PII redaction · citation · schema"),
    ]

    band_top = Inches(2.6)
    band_gap = Inches(0.15)
    band_h = (H - band_top - Inches(0.6) - band_gap * (len(tbs) - 1)) / len(tbs)

    for i, (code, direction, fill, border, tcol, controls) in enumerate(tbs):
        y = band_top + (band_h + band_gap) * i
        # code chip
        chip_w = Inches(2.7)
        add_box(
            slide, Inches(0.35), y, chip_w, band_h,
            border, border, direction,
            title=code, title_size=13, font_size=10,
            title_color="#FFFFFF", text_color="#FFECEB", align_left=True,
        )
        # controls box
        body_left = Inches(0.35) + chip_w + Inches(0.1)
        body_w = W - body_left - Inches(0.35)
        add_box(
            slide, body_left, y, body_w, band_h,
            fill, border, controls,
            title="Controls at this boundary",
            title_size=11, font_size=11,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H)


def build_runtime_flow(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Runtime execution flow (D-R1)")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "One request through the seven decision points: identity, cache, retrieval, guardrails, model, tool, verification.",
        size=12, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    stages = [
        ("1", "Ingress",
         "L0 / L1: AuthN → AuthZ → rate limit → request validation → session",
         "entry"),
        ("2", "Cache lookup",
         "L11: prompt / semantic / response cache — return if hit",
         "cache"),
        ("3", "Context assembly",
         "L2 → L3: agent selection · context builder · memory + retrieval",
         "neutral"),
        ("4", "Retrieval",
         "L4 + L5: memory + hybrid retrieval (vector · BM25 · SQL · graph) → rerank",
         "data"),
        ("5", "Prompt + input guardrails",
         "L3 → L8: prompt builder → input / prompt guardrails",
         "security"),
        ("6", "Model routing",
         "L6: capability · cost · latency · policy → foundation or specialised model",
         "model"),
        ("7", "Tool execution",
         "L7 + L8: tool router → impact-class enforcer → MCP / enterprise API",
         "data"),
        ("8", "Output guardrails + verification",
         "L8 → L9: output guardrails → grounding · citation · factuality · policy · safety",
         "security"),
        ("9", "Format + stream + response",
         "L9 → L1 → L0: response formatter → streaming engine → user",
         "obs"),
    ]

    top = Inches(1.35)
    n = len(stages)
    row_gap = Inches(0.06)
    row_h = (H - top - Inches(0.5) - row_gap * (n - 1)) / n

    for i, (num, name, body, kind) in enumerate(stages):
        y = top + (row_h + row_gap) * i
        fill, border, tcol = PALETTE[kind]
        # circle number
        cw = Inches(0.55)
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.35), y + (row_h - cw) / 2, cw, cw,
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = hex_to_rgb(border)
        circle.line.color.rgb = hex_to_rgb(border)
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = Inches(0.35) + cw + Inches(0.15)
        body_w = W - body_left - Inches(0.35)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, body,
            title=name, title_size=12, font_size=10.5,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H)


def build_llmops(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="L12 · LLMOps & Delivery")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Traditional CI/CD is not enough. LLMOps adds prompt tests, eval gates, model registry, canary + shadow, cost guardrails.",
        size=12, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    # top row: 5 stages pipeline
    stages = [
        ("Git / Repo",
         "Prompts · agents · tools · IaC · eval sets\nversioned as code",
         "neutral"),
        ("CI",
         "Lint · sec-scan · unit · prompt tests\nsafety tests · eval gate",
         "obs"),
        ("Artefact Registry",
         "Signed containers · pinned models\nprompt versions · SBOM",
         "data"),
        ("CD",
         "Blue / Green · Canary · Shadow\nRollback · progressive delivery",
         "model"),
        ("Runtime",
         "Continuous eval · drift · cost cap\nincident + rollback triggers",
         "security"),
    ]
    top = Inches(1.4)
    row_h = Inches(1.4)
    n = len(stages)
    gap = Inches(0.15)
    total_w = W - Inches(0.7)
    box_w = (total_w - gap * (n - 1)) / n
    for i, (title, body, kind) in enumerate(stages):
        x = Inches(0.35) + (box_w + gap) * i
        fill, border, tcol = PALETTE[kind]
        add_box(
            slide, x, top, box_w, row_h,
            fill, border, body,
            title=title, title_size=13, font_size=10,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    # arrows between them
    ay = top + row_h / 2
    for i in range(n - 1):
        x1 = Inches(0.35) + (box_w + gap) * i + box_w
        x2 = x1 + gap
        arrow_between(slide, x1, ay, x2, ay)

    # second half: three side-by-side registries + evaluation
    reg_top = top + row_h + Inches(0.5)
    reg_h = Inches(2.6)
    reg_w = (W - Inches(0.7) - gap * 2) / 3
    regs = [
        ("Prompt lifecycle",
         "Draft → Test → Approved → Prod → Deprecated",
         [
             "Versioned prompt templates",
             "Prompt tests + regression suite",
             "Prompt registry as source of truth",
             "Prompt rollback independent of model",
         ],
         "neutral"),
        ("Model lifecycle",
         "Candidate → Shadow → Canary → Prod → Retired",
         [
             "Model registry (foundation + fine-tuned)",
             "Approved providers only (TB-C)",
             "Cost + latency + safety benchmarks",
             "Shadow traffic for regression",
         ],
         "model"),
        ("Tool lifecycle",
         "Onboarded → Restricted → GA → Deprecated",
         [
             "Tool registry with impact class",
             "Argument guardrails per tool",
             "HITL approval for high-impact",
             "Audit + rate limit per tool",
         ],
         "data"),
    ]
    for i, (title, sub, items, kind) in enumerate(regs):
        x = Inches(0.35) + (reg_w + gap) * i
        fill, border, tcol = PALETTE[kind]
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, reg_top, reg_w, reg_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb(fill)
        shp.line.color.rgb = hex_to_rgb(border)
        shp.line.width = Pt(1.25)
        tf = shp.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.12)
        tf.margin_bottom = Inches(0.12)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.color.rgb = hex_to_rgb(tcol)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.italic = True
        r2.font.name = "Segoe UI"
        r2.font.color.rgb = hex_to_rgb(tcol)

        for item in items:
            pi = tf.add_paragraph()
            pi.alignment = PP_ALIGN.LEFT
            pi.level = 0
            ri = pi.add_run()
            ri.text = "•  " + item
            ri.font.size = Pt(11)
            ri.font.name = "Segoe UI"
            ri.font.color.rgb = hex_to_rgb(tcol)

    slide_footer(slide, W, H)


def build_cloud_mapping(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Cloud implementations")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Same 16 layers. Four independent mappings. Vendor-neutral by construction.",
        size=12, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    clouds = [
        ("Microsoft Azure", "entry", [
            ("L1 Gateway", "APIM · Front Door · Entra ID"),
            ("L2 Orch.",   "Container Apps · Durable Fn"),
            ("L5 Retrieval", "AI Search · Cosmos Gremlin · Fabric"),
            ("L6 Models", "Azure OpenAI · AI Foundry"),
            ("L7 Tools", "APIM · Logic Apps · MCP on ACA"),
            ("L10 Obs.", "App Insights · Azure Monitor"),
            ("L13 Sec.", "Entra · Key Vault · Purview · Defender"),
            ("L14 GRC", "Purview · Azure Policy · Foundry gov."),
        ]),
        ("Amazon AWS", "data", [
            ("L1 Gateway", "CloudFront · WAF · API Gateway · Cognito"),
            ("L2 Orch.",   "ECS / Fargate · Step Functions"),
            ("L5 Retrieval", "OpenSearch · Neptune · Athena"),
            ("L6 Models", "Bedrock · SageMaker"),
            ("L7 Tools", "API Gateway · Lambda · MCP on ECS"),
            ("L10 Obs.", "CloudWatch · X-Ray · OTel Collector"),
            ("L13 Sec.", "IAM · KMS · Secrets Manager · Macie"),
            ("L14 GRC", "Config · Audit Manager · Bedrock Guardrails"),
        ]),
        ("Google Cloud", "model", [
            ("L1 Gateway", "Cloud LB · Cloud Armor · API Gateway · IAP"),
            ("L2 Orch.",   "Cloud Run · Workflows"),
            ("L5 Retrieval", "Vertex AI Search · Spanner Graph"),
            ("L6 Models", "Vertex AI · Gemini · Model Garden"),
            ("L7 Tools", "API Gateway · Cloud Functions · MCP on Run"),
            ("L10 Obs.", "Cloud Ops Suite · OTel Collector"),
            ("L13 Sec.", "IAM · KMS · Secret Manager · SCC"),
            ("L14 GRC", "Assured Workloads · Vertex Model Registry"),
        ]),
        ("Open Source / K8s", "neutral", [
            ("L1 Gateway", "Envoy · Kong · Keycloak · OIDC"),
            ("L2 Orch.",   "K8s + Argo Workflows · LangGraph"),
            ("L5 Retrieval", "Qdrant · Weaviate · pgvector · Neo4j"),
            ("L6 Models", "vLLM · Ollama · TGI · TorchServe"),
            ("L7 Tools", "MCP servers · custom adapters"),
            ("L10 Obs.", "OTel · Prometheus · Grafana · Tempo · Loki"),
            ("L13 Sec.", "SPIFFE · Vault · Falco · Kyverno"),
            ("L14 GRC", "OPA / Gatekeeper · Kyverno · Sigstore"),
        ]),
    ]

    top = Inches(1.4)
    row_h = Inches(5.6)
    n = len(clouds)
    gap = Inches(0.14)
    total_w = W - Inches(0.7)
    box_w = (total_w - gap * (n - 1)) / n

    for i, (cloud_name, kind, items) in enumerate(clouds):
        x = Inches(0.35) + (box_w + gap) * i
        fill, border, tcol = PALETTE[kind]
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, row_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb(fill)
        shp.line.color.rgb = hex_to_rgb(border)
        shp.line.width = Pt(1.25)
        tf = shp.text_frame
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.14)
        tf.margin_bottom = Inches(0.12)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = cloud_name
        r.font.size = Pt(15)
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.color.rgb = hex_to_rgb(tcol)

        for lname, svc in items:
            pl = tf.add_paragraph()
            pl.alignment = PP_ALIGN.LEFT
            rl = pl.add_run()
            rl.text = lname
            rl.font.size = Pt(10.5)
            rl.font.bold = True
            rl.font.name = "Segoe UI"
            rl.font.color.rgb = hex_to_rgb(tcol)

            pv = tf.add_paragraph()
            pv.alignment = PP_ALIGN.LEFT
            rv = pv.add_run()
            rv.text = svc
            rv.font.size = Pt(9.5)
            rv.font.name = "Segoe UI"
            rv.font.color.rgb = hex_to_rgb(tcol)

    slide_footer(slide, W, H)


def build_standards(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Standards & frameworks")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "EASRA maps explicitly to the frameworks your compliance, security and risk teams already recognise.",
        size=12, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    standards = [
        ("NIST AI RMF",  "Govern · Map · Measure · Manage",
         "L14 Governance · L10 Observability", "neutral"),
        ("ISO / IEC 42001", "AI management system",
         "L14 · L12 LLMOps · L15 Value", "neutral"),
        ("EU AI Act",   "Risk categorisation + provider duties",
         "L14 · L9 Verification · L13 Security", "security"),
        ("OWASP LLM Top 10", "LLM-specific application risks",
         "L1 · L6 · L7 · L8 · L9 · L11", "security"),
        ("MITRE ATLAS", "Adversary tactics + techniques for AI",
         "L1 · L8 Guardrails · L13 · L14", "security"),
        ("OpenTelemetry", "Traces / metrics / logs schema",
         "L10 Observability across every layer", "obs"),
        ("Model Context Protocol (MCP)",
         "Standardised tool interface",
         "L7 Tooling & Actions", "data"),
        ("TOGAF",       "Enterprise architecture context",
         "Positions EASRA as an AI-domain reference architecture", "neutral"),
    ]

    cols = 2
    rows = 4
    top = Inches(1.4)
    left = Inches(0.35)
    col_gap = Inches(0.2)
    row_gap = Inches(0.15)
    total_w = W - Inches(0.7)
    col_w = (total_w - col_gap) / cols
    total_h = H - top - Inches(0.5)
    row_h = (total_h - row_gap * (rows - 1)) / rows

    for i, (name, what, mapped, kind) in enumerate(standards):
        col = i // rows
        row = i % rows
        x = left + (col_w + col_gap) * col
        y = top + (row_h + row_gap) * row
        fill, border, tcol = PALETTE[kind]
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, row_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb(fill)
        shp.line.color.rgb = hex_to_rgb(border)
        shp.line.width = Pt(1.25)
        tf = shp.text_frame
        tf.margin_left = Inches(0.14)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.10)
        tf.margin_bottom = Inches(0.10)
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = name
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.name = "Segoe UI"
        r.font.color.rgb = hex_to_rgb(tcol)

        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = what
        r2.font.size = Pt(10.5)
        r2.font.italic = True
        r2.font.name = "Segoe UI"
        r2.font.color.rgb = hex_to_rgb(tcol)

        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        r3 = p3.add_run()
        r3.text = "Mapped in EASRA at:  " + mapped
        r3.font.size = Pt(10)
        r3.font.name = "Segoe UI"
        r3.font.color.rgb = hex_to_rgb(tcol)

    slide_footer(slide, W, H)


def build_deliverables(prs: Presentation):
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Ten deliverables of the standard")

    items = [
        ("1", "Architecture Specification",  "Twelve numbered specs (001–012). The frozen source of truth."),
        ("2", "Handbook",                    "Per-layer chapters and cross-cutting guides."),
        ("3", "Reference Architectures",     "Five views: logical · runtime · deployment · operational · security."),
        ("4", "Cloud Implementations",       "Azure · AWS · GCP · open-source (K8s + CNCF)."),
        ("5", "Benchmarks",                  "Latency · throughput · cost · safety · verification · reliability."),
        ("6", "Architecture Decision Records","Every consequential decision recorded and versioned."),
        ("7", "Security Reference",          "Threats · controls · standards mapping · threat models · red team."),
        ("8", "Verification Reference",      "Classes · checkers · metrics · golden-set methodology."),
        ("9", "LLMOps Guide",                "Delivery · lifecycles · evaluation · cost · incident response."),
        ("10", "Conference Materials",       "Slides · workshops · talks — this deck is one of them."),
    ]

    cols = 2
    rows = 5
    top = Inches(0.95)
    left = Inches(0.35)
    col_gap = Inches(0.2)
    row_gap = Inches(0.1)
    total_w = W - Inches(0.7)
    col_w = (total_w - col_gap) / cols
    total_h = H - top - Inches(0.5)
    row_h = (total_h - row_gap * (rows - 1)) / rows

    for i, (num, name, desc) in enumerate(items):
        col = i // rows
        row = i % rows
        x = left + (col_w + col_gap) * col
        y = top + (row_h + row_gap) * row
        fill, border, tcol = PALETTE["neutral"]

        chip_w = Inches(0.65)
        chip = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + (row_h - chip_w) / 2, chip_w, chip_w)
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb("#3B7DDD")
        chip.line.color.rgb = hex_to_rgb("#3B7DDD")
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.size = Pt(14)
        cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF")
        cr.font.name = "Segoe UI"

        body_left = x + chip_w + Inches(0.15)
        body_w = col_w - chip_w - Inches(0.15)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, desc,
            title=name, title_size=13, font_size=11,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H)


def build_cta(prs: Presentation):
    slide, W, H = slide_bg(prs, color=PALETTE["hero_bg"][0])

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.55), W, Inches(0.15)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(PALETTE["hero_acc"][0])
    bar.line.fill.background()

    add_text(
        slide, Inches(0.9), Inches(0.9), W - Inches(1.8), Inches(0.5),
        "ADOPT · CONTRIBUTE · FORK",
        size=14, bold=True, color="#3B7DDD", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(1.35), W - Inches(1.8), Inches(1.4),
        "EASRA is an open architecture standard.",
        size=42, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(2.55), W - Inches(1.8), Inches(0.7),
        "Vendor-neutral. Standards-aligned. Verifiable. Yours to fork.",
        size=18, italic=True, color="#B8CDEA", align=PP_ALIGN.LEFT,
    )

    ctas = [
        ("★ Star",       "github.com/KarthikeyanDhanakotti/EASRA",         "entry"),
        ("Fork & Adopt", "templates/repository-template/",                 "data"),
        ("Contribute",   "CONTRIBUTING.md  ·  GOVERNANCE.md",              "obs"),
        ("Follow",       "linkedin.com/in/karthikeyan-dhanakotti",         "model"),
    ]
    x = Inches(0.9)
    y = Inches(4.0)
    bw = Inches(2.75)
    bh = Inches(1.4)
    gap = Inches(0.2)
    for label, url, kind in ctas:
        fill, border, _ = PALETTE[kind]
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb("#12335F")
        shp.line.color.rgb = hex_to_rgb(border)
        shp.line.width = Pt(1.5)
        tf = shp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = label
        r1.font.size = Pt(18)
        r1.font.bold = True
        r1.font.color.rgb = hex_to_rgb("#FFFFFF")
        r1.font.name = "Segoe UI"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = url
        r2.font.size = Pt(10)
        r2.font.color.rgb = hex_to_rgb("#8FA9CE")
        r2.font.name = "Segoe UI"
        x += bw + gap

    add_text(
        slide, Inches(0.9), Inches(6.85), W - Inches(1.8), Inches(0.35),
        "Karthikeyan Dhanakotti  ·  Enterprise AI Architect  ·  Independent open-source work",
        size=11, color="#8FA9CE", align=PP_ALIGN.LEFT,
    )


# ---------- v2 builders (post-review) ---------- #

def _connector(slide, kind, x1, y1, x2, y2, color="#5A6470", weight=1.75,
               dashed=False):
    """Add a connector line. `kind` is one of MSO_CONNECTOR members."""
    line = slide.shapes.add_connector(kind, int(x1), int(y1), int(x2), int(y2))
    line.line.color.rgb = hex_to_rgb(color)
    line.line.width = Pt(weight)
    if dashed:
        try:
            from pptx.enum.dml import MSO_LINE_DASH_STYLE
            line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        except Exception:
            pass
    return line


def _arrow_shape(slide, x, y, w, h, direction="right", color="#3B7DDD"):
    """Draw a solid arrow shape (visible arrowhead)."""
    mapping = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down":  MSO_SHAPE.DOWN_ARROW,
        "left":  MSO_SHAPE.LEFT_ARROW,
        "up":    MSO_SHAPE.UP_ARROW,
    }
    a = slide.shapes.add_shape(mapping[direction], int(x), int(y), int(w), int(h))
    a.fill.solid()
    a.fill.fore_color.rgb = hex_to_rgb(color)
    a.line.color.rgb = hex_to_rgb(color)
    a.line.width = Pt(0.25)
    a.text_frame.text = ""
    return a


def build_hero_v2(prs: Presentation):
    """Cover slide — benefit pillars replacing internal-stat badges."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W, H = prs.slide_width, prs.slide_height

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(PALETTE["hero_bg"][0])
    bg.line.fill.background()

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.55), W, Inches(0.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = hex_to_rgb(PALETTE["hero_acc"][0])
    bar.line.fill.background()

    add_text(
        slide, Inches(0.9), Inches(0.85), Inches(11.5), Inches(0.4),
        "OPEN ARCHITECTURE STANDARD  ·  v0.1.0",
        size=13, bold=True, color="#3B7DDD", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(1.25), Inches(11.5), Inches(1.7),
        "EASRA", size=88, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(2.85), Inches(11.5), Inches(0.9),
        "Enterprise AI Systems Reference Architecture",
        size=32, bold=False, color="#FFFFFF", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, Inches(0.9), Inches(3.85), Inches(11.5), Inches(0.65),
        "The open standard for production-grade Enterprise AI.",
        size=20, italic=True, color="#B8CDEA", align=PP_ALIGN.LEFT,
    )

    pillars = [
        ("VENDOR-NEUTRAL",     "No lock-in. Every layer maps to ≥ 2 clouds."),
        ("PRODUCTION-READY",   "Runtime, deployment & scaling — not a slideware framework."),
        ("VERIFICATION-FIRST", "Grounding, citation, factuality baked in — verification ≠ evaluation."),
        ("CLOUD-AGNOSTIC",     "Azure · AWS · GCP · OSS-K8s — same architecture, four impls."),
    ]
    x = Inches(0.9)
    y = Inches(4.85)
    bw = Inches(2.75)
    bh = Inches(1.55)
    gap = Inches(0.15)
    for label, sub in pillars:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, bw, bh)
        shp.fill.solid()
        shp.fill.fore_color.rgb = hex_to_rgb("#12335F")
        shp.line.color.rgb = hex_to_rgb("#3B7DDD")
        shp.line.width = Pt(1.5)
        tf = shp.text_frame
        tf.margin_left = Inches(0.15); tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.15); tf.margin_bottom = Inches(0.12)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.word_wrap = True
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = label
        r1.font.size = Pt(15); r1.font.bold = True
        r1.font.color.rgb = hex_to_rgb("#3B7DDD"); r1.font.name = "Segoe UI"
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(11); r2.font.color.rgb = hex_to_rgb("#E4EEFB")
        r2.font.name = "Segoe UI"
        x += bw + gap

    add_text(
        slide, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.35),
        "Karthikeyan Dhanakotti  ·  github.com/KarthikeyanDhanakotti/EASRA  "
        "·  CC-BY-4.0 (docs) + Apache-2.0 (code)",
        size=11, color="#8FA9CE", align=PP_ALIGN.LEFT,
    )


def build_capability_ladder(prs: Presentation):
    """EASRA as a TOGAF/OpenTelemetry-style progression."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="The EASRA capability ladder")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.55),
        "EASRA is not just a diagram. It is a full stack — like TOGAF, like OpenTelemetry.",
        size=14, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    steps = [
        ("Capability Model",         "What capabilities every Enterprise AI system must have.", "entry"),
        ("Reference Arch.",          "16 layers, 4 trust boundaries, 3 planes — the shared shape.", "neutral"),
        ("Reference Impls.",         "Azure · AWS · GCP · OSS-K8s — same architecture, four impls.", "model"),
        ("Patterns",                 "Named, reusable solutions: RAG, agent, tool-use, HITL, cache-first.", "data"),
        ("Verification",             "Grounding · citation · factuality · policy — separated from evaluation.", "obs"),
        ("Benchmarks",               "Latency · cost · safety · reliability — comparable across impls.", "cache"),
        ("Certification",            "Conformance profiles: what it means to be ‘EASRA-compliant’.", "security"),
    ]

    top = Inches(1.95)
    n = len(steps)
    gap = Inches(0.12)
    total_w = W - Inches(0.7)
    box_w = (total_w - gap * (n - 1)) / n
    box_h = Inches(4.4)

    for i, (name, body, kind) in enumerate(steps):
        x = Inches(0.35) + (box_w + gap) * i
        y = top
        fill, border, tcol = PALETTE[kind]

        # step card FIRST so the numbered chip can sit on top of it
        add_box(
            slide, x, y, box_w, box_h,
            fill, border, body,
            title=name, title_size=13, font_size=11,
            title_color=tcol, text_color=tcol, align_left=True,
        )

        # numbered chip — bigger, higher, drawn AFTER the card so it's visible
        chip_w = Inches(0.65); chip_h = Inches(0.65)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x + box_w / 2 - chip_w / 2), int(y - chip_h + Inches(0.18)),
            chip_w, chip_h,
        )
        chip.fill.solid(); chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb("#FFFFFF")
        chip.line.width = Pt(2)
        ctf = chip.text_frame
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1)
        cr.font.size = Pt(20); cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF"); cr.font.name = "Segoe UI"

        # arrow to next step (horizontal, level with card middle)
        if i < n - 1:
            ax1 = x + box_w
            ay = y + box_h / 2 - Inches(0.14)
            _arrow_shape(slide, ax1 - Inches(0.02), ay, gap + Inches(0.04),
                         Inches(0.28), "right", "#3B7DDD")

    add_text(
        slide, Inches(0.35), H - Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Ladders 1–3 are v0.1.0 today. Ladders 4–7 are the roadmap.",
        size=12, italic=True, color="#0B3D91", align=PP_ALIGN.LEFT,
    )
    slide_footer(slide, W, H)


def build_architecture_v2(prs: Presentation):
    """Production runtime architecture — U-shaped request/response with
    parallel core, event bus sidebar, cross-cutting bottom strip."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Production runtime architecture")

    add_text(
        slide, Inches(0.35), Inches(0.78), W - Inches(0.7), Inches(0.42),
        "Request flows down through the pipeline; response streams back through verification and guardrails. "
        "Async work fans out to an event bus.",
        size=11, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    # main working area
    left_pad = Inches(0.35)
    top_area = Inches(1.28)
    sidebar_w = Inches(2.6)
    main_w = W - left_pad - sidebar_w - Inches(0.4) - Inches(0.35)  # gap + margin
    sidebar_left = left_pad + main_w + Inches(0.4)

    # -------- Row 1: Ingress path (top strip) -------- #
    row1_top = top_area
    row1_h = Inches(0.55)
    ingress = [
        ("User",         "entry"),
        ("CDN / Edge",   "entry"),
        ("WAF · DDoS",   "security"),
        ("API Gateway",  "entry"),
        ("AI Gateway",   "model"),
    ]
    n1 = len(ingress)
    gap1 = Inches(0.12)
    bw1 = (main_w - gap1 * (n1 - 1)) / n1
    for i, (label, kind) in enumerate(ingress):
        x = left_pad + (bw1 + gap1) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, row1_top, bw1, row1_h, fill, border, label,
                font_size=11, text_color=tcol, title_size=11, bold_title=True)
        if i < n1 - 1:
            ax = x + bw1
            ay = row1_top + row1_h / 2 - Inches(0.11)
            _arrow_shape(slide, ax - Inches(0.02), ay, gap1 + Inches(0.04),
                         Inches(0.22), "right", "#3B7DDD")

    # Distinguish API Gateway (HTTP concerns) from AI Gateway (model concerns)
    api_x = left_pad + (bw1 + gap1) * 3
    ai_x  = left_pad + (bw1 + gap1) * 4
    cap_y = row1_top + row1_h + Inches(0.02)
    # Left-aligned so captions sit clear of the centered down-arrow
    add_text(
        slide, api_x + Inches(0.05), cap_y, bw1 - Inches(0.1), Inches(0.28),
        "OAuth · TLS · quota",
        size=8.5, italic=True, color="#3B7DDD", align=PP_ALIGN.LEFT,
    )
    add_text(
        slide, ai_x + Inches(0.05), cap_y, bw1 - Inches(0.1), Inches(0.28),
        "Routing · retries",
        size=8.5, italic=True, color="#8E44AD", align=PP_ALIGN.LEFT,
    )

    # -------- Row 2: Router / Planner (centered wide box) -------- #
    row2_top = row1_top + row1_h + Inches(0.35)
    row2_h = Inches(0.55)
    router_w = main_w * 0.7
    router_x = left_pad + (main_w - router_w) / 2
    fill, border, tcol = PALETTE["neutral"]
    add_box(slide, router_x, row2_top, router_w, row2_h, fill, border,
            "L2 Router · Planner · Agent Selector · Session Context",
            font_size=11, text_color=tcol, title_size=11, bold_title=True)
    # arrow AI Gateway → Router (vertical) — offset to the RIGHT side of the AI
    # Gateway column so the caption "Routing · retries" doesn't hide the shaft.
    x_ag = left_pad + (bw1 + gap1) * 4 + bw1 * 0.72
    _arrow_shape(slide, x_ag - Inches(0.11), row1_top + row1_h + Inches(0.02),
                 Inches(0.22), Inches(0.32), "down", "#3B7DDD")

    # -------- Row 3: AI reasoning core (5 parallel columns) -------- #
    row3_top = row2_top + row2_h + Inches(0.35)
    row3_h = Inches(2.05)
    core = [
        ("L4 Memory",       "Session · long-term\nprofile · episodic",       "data"),
        ("L5 Retrieval",    "Vector · BM25\nSQL · Graph · rerank",           "data"),
        ("L3 Prompt Build", "Template · pack\ncontext · budget",             "neutral"),
        ("L7 Tools · MCP",  "Registry · impact\nclass · guardrails",         "data"),
        ("L6 Model Router", "Foundation · SLM\ncost · latency · policy",     "model"),
    ]
    n3 = len(core)
    gap3 = Inches(0.12)
    bw3 = (main_w - gap3 * (n3 - 1)) / n3
    for i, (t, body, kind) in enumerate(core):
        x = left_pad + (bw3 + gap3) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, row3_top, bw3, row3_h, fill, border, body,
                title=t, title_size=11, font_size=10,
                title_color=tcol, text_color=tcol, align_left=True)
        # request arrows down from Router
        rx = router_x + router_w * (i + 0.5) / n3
        _arrow_shape(slide, rx - Inches(0.09),
                     row2_top + row2_h + Inches(0.02),
                     Inches(0.18), Inches(0.28), "down", "#8A94A0")
        # response arrows back up (dashed teal, offset)
        _connector(slide, MSO_CONNECTOR.STRAIGHT,
                   x + bw3 - Inches(0.15), row3_top,
                   x + bw3 - Inches(0.15), row2_top + row2_h,
                   color="#16A085", weight=1.25, dashed=True)

    # -------- Row 4: Response pipeline (right-to-left: response starts at
    # L9 Verification rightmost, flows LEFT back toward the User) -------- #
    row4_top = row3_top + row3_h + Inches(0.3)
    row4_h = Inches(0.6)
    verify = [
        ("L1 Streaming",     "entry"),
        ("Response Format",  "obs"),
        ("L8 Guardrails",    "security"),
        ("L9 Verification",  "security"),
    ]
    n4 = len(verify)
    gap4 = Inches(0.14)
    bw4 = (main_w - gap4 * (n4 - 1)) / n4
    for i, (label, kind) in enumerate(verify):
        x = left_pad + (bw4 + gap4) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, row4_top, bw4, row4_h, fill, border, label,
                font_size=11, text_color=tcol, title_size=11, bold_title=True)
        # LEFT-pointing arrow from this box back to the previous (response flow)
        if i > 0:
            ax = x - gap4 - Inches(0.02)
            ay = row4_top + row4_h / 2 - Inches(0.11)
            _arrow_shape(slide, ax, ay,
                         gap4 + Inches(0.04), Inches(0.22),
                         "left", "#C0392B")

    # arrow from L6 Model Router (rightmost row 3) down to L9 Verification
    # (rightmost row 4) — column-aligned so the response flow reads naturally.
    verif_cx = left_pad + (bw4 + gap4) * (n4 - 1) + bw4 / 2
    _arrow_shape(slide, verif_cx - Inches(0.11),
                 row3_top + row3_h + Inches(0.02),
                 Inches(0.22), Inches(0.22), "down", "#C0392B")

    # -------- Response back to User (curved-ish left arrow along the bottom) -------- #
    row5_top = row4_top + row4_h + Inches(0.28)
    row5_h = Inches(0.42)
    add_text(
        slide, left_pad, row5_top, main_w, row5_h,
        "← Streaming response back to User  ·  same pipeline in reverse  ·  citations · redaction · verdict",
        size=11, italic=True, color="#0B3D91", align=PP_ALIGN.CENTER,
    )

    # -------- Right sidebar: Event bus + async workers + DLQ -------- #
    sb_top = top_area
    sb_h = row4_top + row4_h - sb_top
    sb_fill, sb_border, sb_tcol = PALETTE["obs"]
    add_box(slide, sidebar_left, sb_top, sidebar_w, Inches(0.5),
            sb_fill, sb_border, "Kafka · Event Hub · Pub/Sub",
            title="Event Bus", title_size=12, font_size=10,
            title_color=sb_tcol, text_color=sb_tcol, align_left=True)

    sb_items = [
        ("Async Workers",  "Long-running · batch\nembedding · re-index"),
        ("DLQ",            "Dead-letter queue\npoison messages"),
        ("Retry Queue",    "Exponential backoff\nidempotent handlers"),
        ("Outbox / Saga",  "Reliable side-effects\nexactly-once semantics"),
    ]
    item_h = Inches(0.88)
    item_gap = Inches(0.1)
    it_top = sb_top + Inches(0.5) + Inches(0.15)
    for i, (t, body) in enumerate(sb_items):
        y = it_top + (item_h + item_gap) * i
        fill, border, tcol = PALETTE["cache"]
        add_box(slide, sidebar_left, y, sidebar_w, item_h,
                fill, border, body,
                title=t, title_size=11, font_size=9.5,
                title_color=tcol, text_color=tcol, align_left=True)

    # dashed L-shape: core (top of row 3) → UP the gutter → into Event Bus.
    # Terminates at the Event Bus card (top of sidebar), not the Retry Queue.
    bend_x = left_pad + main_w + Inches(0.20)
    eb_y   = sb_top + Inches(0.25)  # vertical center of Event Bus card
    _connector(slide, MSO_CONNECTOR.STRAIGHT,
               bend_x, row3_top,
               bend_x, eb_y,
               color="#B58900", weight=1.5, dashed=True)
    _connector(slide, MSO_CONNECTOR.STRAIGHT,
               bend_x, eb_y,
               sidebar_left, eb_y,
               color="#B58900", weight=1.5, dashed=True)

    # -------- Bottom strip: cross-cutting planes -------- #
    strip_top = H - Inches(0.85)
    strip_h = Inches(0.48)
    strip_w = W - Inches(0.7)
    planes = [
        ("L10  Observability Plane",  "OpenTelemetry · traces · metrics · logs · token/cost accounting"),
        ("L13  Security Plane",       "Zero trust · IAM · secrets · encryption · threat detection · policy"),
        ("L15  Cost & Value Plane",   "Token budgets · autoscaling · GPU scheduling · load shedding · backpressure"),
    ]
    n_p = len(planes)
    gap_p = Inches(0.12)
    bwp = (strip_w - gap_p * (n_p - 1)) / n_p
    for i, (t, body) in enumerate(planes):
        x = Inches(0.35) + (bwp + gap_p) * i
        fill, border, tcol = PALETTE["obs"]
        add_box(slide, x, strip_top, bwp, strip_h, fill, border, body,
                title=t, title_size=10.5, font_size=9,
                title_color=tcol, text_color=tcol, align_left=True)

    slide_footer(slide, W, H,
                 text="Solid arrows = request path  ·  Dashed teal = response  "
                      "·  Dashed amber = async / eventing  ·  Bottom strip = cross-cutting planes")


def build_runtime_sequence(prs: Presentation):
    """Runtime execution as a sequence diagram (swimlanes + arrows)."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Runtime sequence  ·  request lifecycle")

    add_text(
        slide, Inches(0.35), Inches(0.82), W - Inches(0.7), Inches(0.42),
        "One request across the actors that actually cooperate at runtime. "
        "Read top-to-bottom; each arrow is a real hop with real latency.",
        size=11, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    actors = [
        ("User",         "entry"),
        ("AI Gateway",   "model"),
        ("Planner",      "neutral"),
        ("Memory",       "data"),
        ("Retrieval",    "data"),
        ("Model",        "model"),
        ("Tools",        "data"),
        ("Verification", "security"),
    ]
    n = len(actors)
    top_bar_y = Inches(1.4)
    top_bar_h = Inches(0.55)
    lane_top = top_bar_y + top_bar_h
    lane_bottom = H - Inches(1.1)
    lane_h = lane_bottom - lane_top

    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    col_w = total_w / n

    # headers + vertical lifelines
    for i, (name, kind) in enumerate(actors):
        cx = left_pad + col_w * i + col_w / 2
        # header pill
        pill_w = col_w - Inches(0.15)
        pill_x = cx - pill_w / 2
        fill, border, tcol = PALETTE[kind]
        add_box(slide, pill_x, top_bar_y, pill_w, top_bar_h,
                fill, border, name,
                font_size=12, text_color=tcol, title_size=12, bold_title=True)
        # dashed lifeline
        _connector(slide, MSO_CONNECTOR.STRAIGHT,
                   cx, lane_top, cx, lane_bottom,
                   color="#B8C1CC", weight=1.0, dashed=True)

    # sequence steps: (from_idx, to_idx, label, y_offset_fraction, color)
    def lane_x(i):
        return int(left_pad + col_w * i + col_w / 2)

    steps = [
        (0, 1, "1  request + auth",                     "#3B7DDD"),
        (1, 2, "2  policy pass · plan",                 "#3B7DDD"),
        (2, 3, "3  fetch session + memory",             "#16A085"),
        (3, 2, "     history + profile",                "#16A085"),
        (2, 4, "4  retrieve (vector · BM25 · graph)",   "#16A085"),
        (4, 2, "     ranked chunks",                    "#16A085"),
        (2, 5, "5  prompt + input guardrails · call",   "#8E44AD"),
        (5, 2, "     draft answer + tool calls",        "#8E44AD"),
        (2, 6, "6  execute tool (MCP · impact class)",  "#B58900"),
        (6, 2, "     tool result",                      "#B58900"),
        (2, 7, "7  verify (ground · cite · policy)",    "#C0392B"),
        (7, 2, "     verdict + citations",              "#C0392B"),
        (2, 1, "8  final response",                     "#3B7DDD"),
        (1, 0, "9  stream response",                    "#3B7DDD"),
    ]

    n_steps = len(steps)
    # distribute y positions across lane height
    label_h = Inches(0.26)
    step_gap = (lane_h - Inches(0.4)) / n_steps
    for k, (a, b, label, color) in enumerate(steps):
        y = int(lane_top + Inches(0.15) + step_gap * k)
        x1 = lane_x(a); x2 = lane_x(b)
        # arrow line + head
        left_end = min(x1, x2); right_end = max(x1, x2)
        _connector(slide, MSO_CONNECTOR.STRAIGHT,
                   left_end, y, right_end, y,
                   color=color, weight=1.75)
        head_w = Inches(0.18); head_h = Inches(0.22)
        if x2 > x1:  # right-pointing
            _arrow_shape(slide, right_end - head_w, y - head_h / 2,
                         head_w, head_h, "right", color)
            tx = left_end + Inches(0.05)
            tw = right_end - left_end - Inches(0.25)
            align = PP_ALIGN.LEFT
        else:        # left-pointing
            _arrow_shape(slide, left_end, y - head_h / 2,
                         head_w, head_h, "left", color)
            tx = left_end + Inches(0.2)
            tw = right_end - left_end - Inches(0.25)
            align = PP_ALIGN.LEFT
        add_text(slide, tx, y - Inches(0.24), tw, label_h,
                 label, size=9.5, color=color, align=align, bold=True)

    slide_footer(slide, W, H,
                 text="Sequence diagram · not a layered stack. Every arrow is a hop, a timeout, "
                      "a retry, an SLO. Solid arrow = call ·  return arrow = response.")


def build_verification_pipeline(prs: Presentation):
    """L9 Verification expanded to a 10-stage pipeline."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="L9 · Verification pipeline")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Verification ≠ evaluation. Every model output passes through this pipeline before it reaches the user.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    # top: input & output pills
    stages = [
        ("Grounding",      "Is every claim supported by retrieved context?"),
        ("Citation",       "Are sources attached, valid, resolvable?"),
        ("Semantic",       "Does the answer address the actual question?"),
        ("Constraint",     "Structural: schema · format · length · language."),
        ("Execution",      "Generated code / SQL executes without harm."),
        ("Policy",         "Enterprise policy · role · jurisdiction · data class."),
        ("Safety",         "Toxicity · self-harm · CSAM · disallowed content."),
        ("Business Rules", "Pricing · discount · escalation · SLA guardrails."),
        ("Confidence",     "Calibrated uncertainty score · abstain threshold."),
        ("Approval",       "HITL for high-impact · async review queue."),
    ]

    # two rows of 5
    top1 = Inches(1.55)
    row_gap = Inches(0.55)
    row_h = Inches(2.05)
    cols = 5
    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    gap = Inches(0.14)
    bw = (total_w - gap * (cols - 1)) / cols

    for i, (name, body) in enumerate(stages):
        row = i // cols
        col = i % cols
        x = left_pad + (bw + gap) * col
        y = top1 + (row_h + row_gap) * row
        fill, border, tcol = PALETTE["security"]

        # draw card FIRST so chip can render on top
        add_box(slide, x, y, bw, row_h, fill, border, body,
                title=name, title_size=13, font_size=10.5,
                title_color=tcol, text_color=tcol, align_left=True)

        # numbered chip on top-left, drawn on top of the card
        two_digit = (i + 1) >= 10
        chip_w = Inches(0.78 if two_digit else 0.60)
        chip_h = Inches(0.60)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(x + Inches(0.12)),
            int(y - chip_h + Inches(0.18)),
            chip_w, chip_h,
        )
        chip.fill.solid()
        chip.fill.fore_color.rgb = hex_to_rgb("#C0392B")
        chip.line.color.rgb = hex_to_rgb("#FFFFFF")
        chip.line.width = Pt(2)
        ctf = chip.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ctf.margin_left = Emu(0); ctf.margin_right = Emu(0)
        ctf.margin_top = Emu(0); ctf.margin_bottom = Emu(0)
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = str(i + 1)
        cr.font.size = Pt(18); cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF"); cr.font.name = "Segoe UI"

        # right arrow to next in the same row
        if col < cols - 1:
            ax = x + bw
            ay = y + row_h / 2 - Inches(0.11)
            _arrow_shape(slide, ax - Inches(0.02), ay,
                         gap + Inches(0.04), Inches(0.22),
                         "right", "#C0392B")

    # curved-ish arrow row1-end → row2-start (drawn as a down arrow under row1-last)
    last1_x = left_pad + (bw + gap) * (cols - 1) + bw / 2
    first2_x = left_pad + bw / 2
    mid_y = top1 + row_h + Inches(0.15)
    # down from row1 last
    _arrow_shape(slide, last1_x - Inches(0.11), mid_y,
                 Inches(0.22), Inches(0.25), "down", "#C0392B")
    # horizontal line across
    _connector(slide, MSO_CONNECTOR.STRAIGHT,
               last1_x, mid_y + Inches(0.25),
               first2_x, mid_y + Inches(0.25),
               color="#C0392B", weight=1.5)
    # left arrow head down into row 2 first
    _arrow_shape(slide, first2_x - Inches(0.11), mid_y + Inches(0.25),
                 Inches(0.22), Inches(0.25), "down", "#C0392B")

    slide_footer(slide, W, H,
                 text="Any stage may VETO · abstain · escalate to HITL · downgrade confidence · request re-generation.")


def build_platform_services(prs: Presentation):
    """Platform services sitting between AI layer and infrastructure."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Platform services  ·  the middle plane")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.5),
        "Between the AI reasoning layer and cloud infrastructure sits a platform layer of "
        "shared, stateful services. Enterprise AI stands or falls on this middle.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    services = [
        ("L11 Cache Manager",   "Prompt · semantic · embedding\nmemory · response caches",  "cache"),
        ("Kafka / Event Bus",   "Async pipelines · streaming events",       "obs"),
        ("Scheduler",           "Cron · workflows · DAGs · Airflow / Argo", "neutral"),
        ("Configuration",       "App Config · Consul · git-versioned",      "neutral"),
        ("Secrets",             "Key Vault · Vault · KMS · rotation",       "security"),
        ("Feature Flags",       "Progressive delivery · kill-switches",     "obs"),
        ("Workflow Registry",   "Named agents · SOPs · replayable runs",    "data"),
        ("Prompt Registry",     "Versioned prompts · A/B · rollback",       "data"),
        ("Policy Engine",       "OPA · Cedar · content + tool policies",    "security"),
        ("Artifact Registry",   "Signed containers · SBOM · attestations",  "neutral"),
        ("Model Registry",      "Foundation · fine-tuned · shadow · prod",  "model"),
        ("Tool Registry",       "MCP servers · impact class · rate limit",  "data"),
    ]

    cols = 4
    rows = 3
    top = Inches(1.55)
    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    gap_x = Inches(0.14)
    gap_y = Inches(0.14)
    bw = (total_w - gap_x * (cols - 1)) / cols
    total_h = H - top - Inches(0.5)
    bh = (total_h - gap_y * (rows - 1)) / rows

    for i, (name, body, kind) in enumerate(services):
        col = i % cols
        row = i // cols
        x = left_pad + (bw + gap_x) * col
        y = top + (bh + gap_y) * row
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, y, bw, bh, fill, border, body,
                title=name, title_size=13, font_size=10.5,
                title_color=tcol, text_color=tcol, align_left=True)

    slide_footer(slide, W, H)


def build_deployment_view(prs: Presentation):
    """Physical deployment view (Azure-anchored, notes for AWS/GCP/OSS)."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Deployment view  ·  physical architecture")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.45),
        "Logical architecture ≠ deployment architecture. Here is what actually runs, autoscales and pages someone at 3 AM.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    # top row: edge
    edge_top = Inches(1.45)
    edge_h = Inches(0.55)
    edge = [
        ("Internet",        "entry"),
        ("Front Door / CDN",   "entry"),
        ("WAF · DDoS · Bot",   "security"),
        ("Private Endpoint", "security"),
    ]
    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    gap = Inches(0.14)
    n_e = len(edge)
    bwe = (total_w - gap * (n_e - 1)) / n_e
    for i, (label, kind) in enumerate(edge):
        x = left_pad + (bwe + gap) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, edge_top, bwe, edge_h, fill, border, label,
                font_size=11, text_color=tcol, title_size=11, bold_title=True)
        if i < n_e - 1:
            ax = x + bwe
            ay = edge_top + edge_h / 2 - Inches(0.11)
            _arrow_shape(slide, ax - Inches(0.02), ay,
                         gap + Inches(0.04), Inches(0.22),
                         "right", "#3B7DDD")

    # middle: compute cluster (AKS / ECS / GKE / K8s)
    mid_top = edge_top + edge_h + Inches(0.35)
    mid_h = Inches(2.4)
    fill, border, tcol = PALETTE["neutral"]

    # container box for the cluster
    cluster = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left_pad, mid_top, total_w, mid_h,
    )
    cluster.fill.solid()
    cluster.fill.fore_color.rgb = hex_to_rgb("#F7F9FC")
    cluster.line.color.rgb = hex_to_rgb("#2C3E50")
    cluster.line.width = Pt(1.5)
    tf = cluster.text_frame
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.15); tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP; tf.word_wrap = True
    tp = tf.paragraphs[0]; tp.alignment = PP_ALIGN.LEFT
    tr = tp.add_run()
    tr.text = "AKS  ·  Container Apps  ·  (or  ECS · GKE · K8s)   —   HPA + KEDA + GPU scheduler + PodDisruptionBudget"
    tr.font.size = Pt(12); tr.font.bold = True
    tr.font.color.rgb = hex_to_rgb("#1B2A3A"); tr.font.name = "Segoe UI"

    # inner pods
    inner_top = mid_top + Inches(0.55)
    inner_h = mid_h - Inches(0.75)
    pods = [
        ("Gateway pods",     "APIM sidecar\n5–50 replicas\nHPA on RPS",       "entry"),
        ("Orchestrator pods","Container Apps\n3–30 replicas\nHPA on queue",   "neutral"),
        ("Agent workers",    "GPU nodepool\n2–20 replicas\nKEDA on Kafka",    "model"),
        ("Tool executors",   "Sandboxed\n5–30 replicas\nrate limited",        "data"),
        ("Verifier pods",    "CPU nodepool\n3–15 replicas\nHPA on backlog",   "security"),
    ]
    n_p = len(pods)
    gap_p = Inches(0.12)
    bwp = (total_w - Inches(0.3) - gap_p * (n_p - 1)) / n_p
    for i, (t, body, kind) in enumerate(pods):
        x = left_pad + Inches(0.15) + (bwp + gap_p) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, inner_top, bwp, inner_h, fill, border, body,
                title=t, title_size=11, font_size=9.5,
                title_color=tcol, text_color=tcol, align_left=True)

    # bottom row: managed data & AI services + observability sink
    bot_top = mid_top + mid_h + Inches(0.3)
    bot_h = Inches(1.4)
    bot_services = [
        ("Redis Enterprise", "cache + rate-limit\nzone-redundant",      "cache"),
        ("Qdrant / AI Search","vector index\nreplicated",               "data"),
        ("Cosmos DB",        "session + memory\nmulti-region",          "data"),
        ("Azure OpenAI",     "AOAI · Foundry\nprovisioned + PAYG",      "model"),
        ("Blob Storage",     "docs · attachments\nlifecycle rules",     "neutral"),
        ("OpenTelemetry",    "Collector →\nApp Insights / Grafana",     "obs"),
    ]
    n_b = len(bot_services)
    bwb = (total_w - gap * (n_b - 1)) / n_b
    for i, (t, body, kind) in enumerate(bot_services):
        x = left_pad + (bwb + gap) * i
        fill, border, tcol = PALETTE[kind]
        add_box(slide, x, bot_top, bwb, bot_h, fill, border, body,
                title=t, title_size=11, font_size=9.5,
                title_color=tcol, text_color=tcol, align_left=True)

    slide_footer(slide, W, H,
                 text="Same shape on AWS (EKS + ElastiCache + OpenSearch + Bedrock + S3 + CloudWatch) "
                      "and GCP (GKE + Memorystore + Vertex + GCS + Cloud Ops).")


# ---------- main ---------- #

def build_three_planes(prs: Presentation):
    """Control / Data / Management planes — Kubernetes-style separation."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Control · Data · Management planes")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.5),
        "Enterprise AI has three planes — just like Kubernetes. "
        "Confusing them is how governance, cost and reliability go wrong.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    planes = [
        ("Control Plane",
         "Decides · governs · configures",
         "neutral",
         [
             "Policy Engine (OPA · Cedar)",
             "Configuration store",
             "Model Registry",
             "Prompt Registry",
             "Workflow Registry",
             "Tool Registry",
             "Feature Flags",
             "Model routing rules",
             "Rate-limit + quota rules",
         ]),
        ("Data Plane",
         "Runs · executes · serves traffic",
         "data",
         [
             "Prompts & context assembly",
             "Embeddings & retrieval",
             "Model inference calls",
             "Tool executions (MCP)",
             "Verification checkers",
             "Guardrail evaluators",
             "Streaming responses",
             "Cache reads / writes",
             "Event bus messages",
         ]),
        ("Management Plane",
         "Observes · pays · improves",
         "obs",
         [
             "CI / CD pipelines",
             "OpenTelemetry traces & metrics",
             "Cost dashboards & budgets",
             "Alerting & on-call",
             "Governance & audit logs",
             "Compliance evidence",
             "LLMOps eval & benchmarks",
             "Incident response",
             "Capacity + quota planning",
         ]),
    ]

    top = Inches(1.55)
    n = len(planes)
    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    gap = Inches(0.2)
    box_w = (total_w - gap * (n - 1)) / n
    box_h = H - top - Inches(0.55)

    for i, (title, sub, kind, items) in enumerate(planes):
        x = left_pad + (box_w + gap) * i
        fill, border, tcol = PALETTE[kind]
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        shp.fill.solid(); shp.fill.fore_color.rgb = hex_to_rgb(fill)
        shp.line.color.rgb = hex_to_rgb(border); shp.line.width = Pt(1.5)
        tf = shp.text_frame
        tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.14)
        tf.margin_top = Inches(0.18); tf.margin_bottom = Inches(0.14)
        tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.TOP

        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run(); r1.text = title
        r1.font.size = Pt(18); r1.font.bold = True
        r1.font.color.rgb = hex_to_rgb(tcol); r1.font.name = "Segoe UI"

        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(11); r2.font.italic = True
        r2.font.color.rgb = hex_to_rgb(tcol); r2.font.name = "Segoe UI"

        # spacer
        ps = tf.add_paragraph()
        rs = ps.add_run(); rs.text = " "
        rs.font.size = Pt(4)

        for item in items:
            pi = tf.add_paragraph(); pi.alignment = PP_ALIGN.LEFT
            ri = pi.add_run(); ri.text = "•  " + item
            ri.font.size = Pt(11); ri.font.name = "Segoe UI"
            ri.font.color.rgb = hex_to_rgb(tcol)

    slide_footer(slide, W, H)


def build_adrs(prs: Presentation):
    """Architecture Decision Records — every big choice is written down."""
    slide, W, H = slide_bg(prs)
    title_bar(slide, W, subtitle_right="Architecture Decision Records (ADRs)")

    add_text(
        slide, Inches(0.35), Inches(0.85), W - Inches(0.7), Inches(0.55),
        "Every consequential architectural choice in EASRA is recorded as an ADR. "
        "Not a deliverable — a first-class artefact of the standard.",
        size=13, italic=True, color="#5A6470", align=PP_ALIGN.LEFT,
    )

    adrs = [
        ("ADR-001", "Verification is separate from Evaluation",
         "Runtime verifiers vs. offline eval — different lifecycles, different owners.", "security"),
        ("ADR-002", "Memory is externalised",
         "Compute stays stateless. Sessions, profiles, episodic memory live outside.", "data"),
        ("ADR-003", "Prompt Registry is mandatory",
         "Prompts are versioned code, not literals — enables A/B, rollback, audit.", "data"),
        ("ADR-004", "AI Gateway is separate from API Gateway",
         "Model routing, token budgets, retries, circuit breakers ≠ HTTP concerns.", "model"),
        ("ADR-005", "Tool Registry with impact class is mandatory",
         "Every tool has an impact class; HITL for high-impact; audit for all.", "data"),
        ("ADR-006", "Cache Manager is centralised",
         "One component for prompt · semantic · embedding · memory · response caches.", "cache"),
        ("ADR-007", "Async work uses an Event Bus, not sync calls",
         "Batch, embedding, re-index, long-running agents — Kafka / Event Hub.", "obs"),
        ("ADR-008", "Control / Data / Management planes are separated",
         "Governance, runtime and operations scale independently.", "neutral"),
        ("ADR-009", "Standards are mapped, not replaced",
         "NIST · ISO 42001 · EU AI Act · OWASP LLM · MITRE ATLAS — complements, doesn’t compete.", "neutral"),
        ("ADR-010", "Dual license: CC-BY-4.0 docs · Apache-2.0 code",
         "Docs freely adaptable; code freely usable in commercial products.", "entry"),
    ]

    cols = 2
    rows = 5
    top = Inches(1.55)
    left_pad = Inches(0.35)
    total_w = W - Inches(0.7)
    col_gap = Inches(0.2)
    row_gap = Inches(0.12)
    col_w = (total_w - col_gap) / cols
    total_h = H - top - Inches(0.5)
    row_h = (total_h - row_gap * (rows - 1)) / rows

    for i, (code, name, body, kind) in enumerate(adrs):
        col = i // rows
        row = i % rows
        x = left_pad + (col_w + col_gap) * col
        y = top + (row_h + row_gap) * row
        fill, border, tcol = PALETTE[kind]

        # ADR code chip on the left
        chip_w = Inches(1.05)
        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y, chip_w, row_h)
        chip.fill.solid(); chip.fill.fore_color.rgb = hex_to_rgb(border)
        chip.line.color.rgb = hex_to_rgb(border)
        ctf = chip.text_frame
        ctf.margin_left = Inches(0.05); ctf.margin_right = Inches(0.05)
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE; ctf.word_wrap = True
        cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run(); cr.text = code
        cr.font.size = Pt(13); cr.font.bold = True
        cr.font.color.rgb = hex_to_rgb("#FFFFFF"); cr.font.name = "Segoe UI"

        body_left = x + chip_w + Inches(0.1)
        body_w = col_w - chip_w - Inches(0.1)
        add_box(
            slide, body_left, y, body_w, row_h,
            fill, border, body,
            title=name, title_size=12, font_size=10.5,
            title_color=tcol, text_color=tcol, align_left=True,
        )

    slide_footer(slide, W, H,
                 text="ADRs live under docs/adr/. Each one has context · decision · consequences · status · date.")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_hero_v2(prs)                # 1   Cover (benefits, not stats)
    build_problem(prs)                # 2   Problem
    build_what_is(prs)                # 3   What is EASRA?
    build_capability_ladder(prs)      # 4   Capability ladder
    build_principles(prs)             # 5   Ten design principles
    build_layer_index(prs)            # 6   Sixteen layers
    build_architecture_v2(prs)        # 7   Production runtime arch
    build_runtime_sequence(prs)       # 8   Sequence diagram
    build_trust_boundaries(prs)       # 9   Four trust boundaries
    build_verification_pipeline(prs)  # 10  Verification pipeline
    build_three_planes(prs)           # 11  Control / Data / Mgmt planes [NEW]
    build_platform_services(prs)      # 12  Platform services
    build_deployment_view(prs)        # 13  Deployment view
    build_llmops(prs)                 # 14  LLMOps & delivery
    build_cloud_mapping(prs)          # 15  Cloud implementations
    build_standards(prs)              # 16  Standards mapping
    build_adrs(prs)                   # 17  Architecture Decision Records [NEW]
    build_cta(prs)                    # 18  CTA

    out = Path(__file__).parent / "EASRA-LinkedIn-Carousel.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB, {len(list(prs.slides))} slides)")


if __name__ == "__main__":
    main()
